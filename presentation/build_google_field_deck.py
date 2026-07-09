from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parent
SRC = ROOT / "google_field_source.pptx"
OUT = ROOT / "google_field_augmented.pptx"


def delete_slide(prs: Presentation, idx: int) -> None:
    slide_id = prs.slides._sldIdLst[idx]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[idx]


def text_shapes(slide):
    shapes = []
    for shape in slide.shapes:
        txt = getattr(shape, "text", None)
        if txt is not None:
            shapes.append(shape)
    return shapes


def set_texts(slide, texts):
    shapes = text_shapes(slide)
    if len(shapes) < len(texts):
        raise ValueError(f"slide has {len(shapes)} text shapes, need {len(texts)}")
    for shape, text in zip(shapes, texts):
        shape.text = text


def add_section_header(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if 0 in placeholders:
        placeholders[0].text = title
    if 1 in placeholders:
        placeholders[1].text = subtitle
    return slide


def add_title_body(prs: Presentation, title: str, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if 0 in placeholders:
        placeholders[0].text = title
    body = placeholders.get(1)
    if body is None:
        return slide
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
    return slide


def main():
    prs = Presentation(str(SRC))

    # Keep only the meaningful topic slides from the source deck.
    for idx in range(len(prs.slides) - 1, 15, -1):
        delete_slide(prs, idx)

    # Rewrite the existing useful slides with cleaner content.
    set_texts(
        prs.slides[2],
        [
            "Style Transfer",
            "What is style transfer?",
            "At a high level, style transfer means taking the content of one piece and combining it with the style of another.\n\n"
            "In music, that means preserving the underlying musical identity while changing timbre, instrumentation, texture, articulation, or production character.",
        ],
    )
    set_texts(
        prs.slides[5],
        [
            "Style vs Content",
            "Content: melody, rhythm, notes being played, harmonic path, phrase structure\n\n"
            "Style: timbre, instrumentation, texture, dynamics, articulation, expressive details",
        ],
    )
    set_texts(
        prs.slides[6],
        [
            "The Challenge",
            "Why music style transfer is harder than it first seems",
            "A single recording contains melody, rhythm, harmony, timbre, texture, instrumentation, room acoustics, mixing decisions, and long-range form all at once.\n\n"
            "Because these factors are entangled, changing one layer often damages another. That is why weak systems often produce a “coat of paint” result rather than a true re-authoring.",
        ],
    )
    set_texts(
        prs.slides[7],
        [
            "Audio Representations",
            "Audio can be represented as:\n\n"
            "MIDI\n"
            "Spectrograms and log-mel features\n"
            "Raw waveforms\n"
            "Learned latent representations",
        ],
    )
    set_texts(
        prs.slides[8],
        [
            "Audio Representations",
            "MIDI\n\n"
            "Stores note events and controller messages rather than final sound\n"
            "Small file size\n"
            "Easy manipulation\n"
            "But a lot of timbral realism is missing",
        ],
    )
    set_texts(
        prs.slides[9],
        [
            "Audio Representations",
            "Raw waveforms\n\n"
            "Highest fidelity representation\n"
            "X-axis is time, Y-axis is air-pressure displacement\n"
            "Wave shape carries timbre and local acoustic detail\n"
            "But it is a very dense space for semantic editing",
        ],
    )
    set_texts(
        prs.slides[10],
        [
            "Audio Representations",
            "Spectrograms\n\n"
            "Commonly used in style-transfer architectures\n"
            "X-axis is time, Y-axis is frequency, color intensity reflects amplitude\n"
            "Computed from waveforms using the Short-Time Fourier Transform\n"
            "They expose musical structure more clearly than raw waveforms",
        ],
    )
    set_texts(
        prs.slides[11],
        [
            "Short Time Fourier Transform",
            "Fourier Transform vs STFT\n\n"
            "The Fourier Transform tells us what frequencies are present overall.\n"
            "The STFT applies the transform to short overlapping windows, giving a time-frequency representation of the signal.",
        ],
    )
    set_texts(
        prs.slides[13],
        [
            "WaveNet",
            "Introduced by DeepMind in 2016\n"
            "Autoregressive model with dilated causal convolutions\n"
            "A major milestone in realistic neural audio generation\n"
            "But it generates one sample at a time, so inference is slow and expensive",
        ],
    )
    set_texts(
        prs.slides[14],
        [
            "Background About GANs\n\n"
            "Generative Adversarial Networks are built around two networks:\n"
            "Generator: produces fake samples\n"
            "Discriminator: distinguishes real samples from fake ones\n\n"
            "They became popular because they often produced sharper outputs than standard autoencoder-style models.",
        ],
    )
    set_texts(
        prs.slides[15],
        [
            "WaveGAN",
            "One of the first successful GAN-based raw-audio generators\n"
            "Adapts image-GAN ideas to one-dimensional audio\n"
            "Introduces phase shuffle to reduce trivial phase artifacts\n\n"
            "Important historically, but limited by phase precession and short fixed-length output",
        ],
    )

    # Add the missing field slides.
    add_title_body(
        prs,
        "GANSynth",
        [
            "Instead of generating raw waveforms directly, GANSynth generates log-magnitude spectrograms and instantaneous-frequency spectrograms.",
            "Instantaneous frequency is more stable than raw phase under alignment shifts.",
            "Magnitude + IF are converted back into waveform audio through phase reconstruction and inverse STFT.",
            "Field lesson: representation choice matters.",
        ],
    )
    add_title_body(
        prs,
        "Why Early Audio Generation Was Not Enough",
        [
            "WaveNet, WaveGAN, and GANSynth showed that neural audio generation was possible.",
            "But they did not explicitly guarantee that the original musical content would be preserved.",
            "They were strong historical steps in audio generation, but weak direct solutions to music style transfer.",
        ],
    )
    add_section_header(prs, "Latent Structured Audio and Style Transfer")
    add_title_body(
        prs,
        "Variational Autoencoders",
        [
            "VAEs are attractive because they expose a latent space where musical properties might be manipulated.",
            "Encoder compresses the input into latent variable z; decoder reconstructs from z.",
            "Naive audio VAEs struggle with posterior collapse, weak global bottlenecks, and raw sequential complexity.",
            "Still, the latent-space idea remains very valuable for style transfer.",
        ],
    )
    add_title_body(
        prs,
        "MoVE",
        [
            "MoVE stands for Modulated Variational auto-Encoders for many-to-many musical timbre transfer.",
            "It was built specifically for musical timbre transfer rather than generic audio generation.",
            "Uses FiLM conditioning and an MMD objective to support many-to-many transfer in a single multi-domain model.",
            "Very relevant conceptually, but still closer to timbre transfer than full genre remastering.",
        ],
    )
    add_title_body(
        prs,
        "RAVE",
        [
            "RAVE makes high-quality latent audio practical.",
            "Uses spectral losses instead of only raw sample-wise reconstruction.",
            "Preserves temporal structure with sequences of latent codes rather than one weak global vector.",
            "Improves decoder realism with adversarial sharpening and supports music-quality audio generation.",
        ],
    )
    add_title_body(
        prs,
        "AudioLM",
        [
            "AudioLM is not a style-transfer model directly, but it addresses long-form coherence.",
            "Treats audio generation like language modeling over discrete audio tokens.",
            "Separates tokens for long-range structure from tokens for fine acoustic detail.",
            "Important because long-range stability is one of the hardest downstream problems in genre remastering.",
        ],
    )
    add_section_header(prs, "Modern Controllable Synthesis and Style Transfer")
    add_title_body(
        prs,
        "BigVGAN and FiLM",
        [
            "BigVGAN is a universal neural vocoder: it turns intermediate audio representations into more realistic waveform audio.",
            "Its role is realism rather than transfer logic itself.",
            "FiLM is a conditioning mechanism that lets a style vector scale and shift hidden features throughout the network.",
            "Its role is style-control logic rather than final waveform synthesis.",
        ],
    )
    add_title_body(
        prs,
        "DDPM",
        [
            "Denoising Diffusion Probabilistic Models gradually destroy a sample with noise and learn the reverse denoising process.",
            "Diffusion matters because generation happens by gradual refinement rather than one-shot synthesis.",
            "That usually gives more stability and more room for control than many older model families.",
        ],
    )
    add_title_body(
        prs,
        "Classifier-Free Guidance and SDEdit",
        [
            "Classifier-Free Guidance adds a practical control knob for how strongly the generation should follow the target condition.",
            "SDEdit starts from a noised version of an existing source instead of pure noise.",
            "Together they make the content-versus-style tradeoff much more explicit.",
        ],
    )
    add_title_body(
        prs,
        "AudioLDM and Why Diffusion Is the Modern Direction",
        [
            "AudioLDM shows that high-quality controllable generation can happen in learned latent audio spaces.",
            "This makes diffusion much more relevant to editing and transformation, not only text-to-audio generation.",
            "Diffusion is central because it combines progressive generation, conditioning, guidance, and source anchoring.",
        ],
    )
    add_title_body(
        prs,
        "Overview of Field Trends",
        [
            "1. Disentanglement",
            "2. Controllable latent audio",
            "3. Diffusion-based re-authoring",
            "4. Long-form coherence",
            "The field moved from raw generation toward structured representations and explicit control.",
        ],
    )
    add_title_body(
        prs,
        "Conclusion",
        [
            "Music style transfer is not just an audio-generation problem.",
            "It is also a structure-preservation, representation, conditioning, synthesis, and coherence problem.",
            "The field is moving from surface transfer toward structure-aware genre reconstruction.",
        ],
    )

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
