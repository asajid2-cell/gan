from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib import rcParams
import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.colors import Color
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.utils import ImageReader
from reportlab.platypus import Paragraph
from reportlab.pdfgen import canvas


REPO_ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = REPO_ROOT / "presentation"
ASSETS_DIR = PRESENTATION_DIR / "assets"
MEDIA_DIR = PRESENTATION_DIR / "media"
GENERATED_DIR = PRESENTATION_DIR / "generated"
FORMULA_DIR = GENERATED_DIR / "formulas"

PPTX_OUT = PRESENTATION_DIR / "dggr_lecture_deck.pptx"
PDF_OUT = PRESENTATION_DIR / "dggr_lecture_deck.pdf"


COLORS = {
    "ink": (31, 41, 55),
    "muted": (99, 115, 129),
    "sand": (247, 244, 236),
    "paper": (255, 253, 248),
    "mist": (232, 241, 239),
    "peach": (246, 228, 216),
    "slate": (221, 228, 236),
    "gold": (194, 125, 56),
    "teal": (20, 122, 115),
    "rose": (200, 92, 92),
    "night": (15, 23, 42),
}


def rgb(name: str) -> RGBColor:
    return RGBColor(*COLORS[name])


def pdf_color(name: str) -> Color:
    r, g, b = COLORS[name]
    return Color(r / 255.0, g / 255.0, b / 255.0)


def ensure_dirs() -> None:
    GENERATED_DIR.mkdir(exist_ok=True)
    FORMULA_DIR.mkdir(exist_ok=True)


def next_available_path(path: Path) -> Path:
    if not path.exists():
        return path
    for idx in range(2, 20):
        candidate = path.with_name(f"{path.stem}_v{idx}{path.suffix}")
        if not candidate.exists():
            return candidate
    return path.with_name(f"{path.stem}_rebuilt{path.suffix}")


def make_formula_image(filename: str, lines: list[str]) -> Path:
    out = FORMULA_DIR / filename
    if out.exists():
        return out

    rcParams["mathtext.fontset"] = "dejavusans"
    fig = plt.figure(figsize=(10, 1.8), dpi=200)
    fig.patch.set_facecolor("#FFFDF8")
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    y = 0.72
    for line in lines:
        ax.text(0.03, y, line, fontsize=18, color="#1F2937", va="center", ha="left")
        y -= 0.34
    fig.savefig(out, dpi=200, facecolor=fig.get_facecolor(), bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)
    return out


def formula_assets() -> dict[str, Path]:
    return {
        "problem": make_formula_image(
            "problem.png",
            [
                r"$x_{\mathrm{src}} \rightarrow \hat{x}_{\mathrm{tgt}}$",
                r"$\mathrm{content}(\hat{x}_{\mathrm{tgt}})\approx \mathrm{content}(x_{\mathrm{src}}),\ \hat{x}_{\mathrm{tgt}} \in \mathcal{M}_{\mathrm{genre=tgt}}$",
            ],
        ),
        "disentangle": make_formula_image(
            "disentangle.png",
            [
                r"$E(x)\rightarrow (z_c,z_s)$",
                r"$\mathcal{L}=\lambda_c\mathcal{L}_{content}+\lambda_s\mathcal{L}_{style}+\lambda_a\mathcal{L}_{adv}$",
            ],
        ),
        "target160": make_formula_image(
            "target160.png",
            [
                r"$\mathrm{target160}=\mathrm{normalize}([2.0\,z_s\ \Vert\ 1.0\,d_{32}])$",
                r"$d_{32}=[\mu_{1:16},\sigma_{1:16}],\quad \mu_g=\frac{1}{|S_g|}\sum_{x\in S_g} v(x)$",
            ],
        ),
        "film": make_formula_image(
            "film.png",
            [
                r"$\mathrm{FiLM}(F;c)=\gamma(c)\odot F+\beta(c)$",
                r"$\mathrm{AdaIN}(F;s)=(1+\gamma(s))\odot \mathrm{IN}(F)+\beta(s)$",
            ],
        ),
        "diffusion": make_formula_image(
            "diffusion.png",
            [
                r"$q(x_t|x_0)=\mathcal{N}(\sqrt{\bar{\alpha}_t}\,x_0,\ (1-\bar{\alpha}_t)I)$",
                r"$\hat{\epsilon}_{cfg}=\epsilon_\theta(x_t,t,\varnothing)+w[\epsilon_\theta(x_t,t,c)-\epsilon_\theta(x_t,t,\varnothing)]$",
            ],
        ),
        "prefix": make_formula_image(
            "prefix.png",
            [
                r"$x_t^{(k)}=\alpha(t)x_0^{(k)}+\sigma(t)\epsilon$",
                r"$\mathrm{prefix}(x_t^{(k)}) \leftarrow \mathrm{tail}(x_t^{(k-1)})$",
            ],
        ),
        "codec_loss": make_formula_image(
            "codec_loss.png",
            [
                r"$\mathcal{L}_{codec}=\lambda_{adv}\mathcal{L}_{GAN}+\lambda_{fm}\mathcal{L}_{feat}+\lambda_1\|q_{hat}-q^*\|_1+\lambda_{mrstft}\mathcal{L}_{MRSTFT}$",
                r"$\qquad\qquad+\lambda_c\mathcal{L}_{content}+\lambda_s\mathcal{L}_{style}+\lambda_p\mathcal{L}_{push}$",
            ],
        ),
    }


def slide_specs(formulas: dict[str, Path]) -> list[dict]:
    return [
        {
            "kind": "title",
            "title": "Music Style Transfer and Deep\nGenerative Genre Remastering",
            "subtitle": "50-minute lecture + DGGR case study",
            "body": [
                "Part I: general lecture on the field, model families, formulas, and evaluation.",
                "Part II: our repo-backed system, what we implemented, and what the results actually say.",
            ],
        },
        {
            "kind": "split",
            "title": "Talk Map",
            "left_title": "Lecture half",
            "left": [
                "Define content, style, and genre before discussing our project.",
                "Cover representations, disentanglement, conditioning, codecs, diffusion, and long-form coherence.",
                "Explain what a good evaluation protocol must measure.",
            ],
            "right_title": "Case-study half",
            "right": [
                "Show DGGR as one concrete architecture built from those ideas.",
                "Walk lab by lab through architecture, metrics, decisions, and remaining failures.",
                "Close with audio-demo guidance, creative elements, and Lab 5.",
            ],
        },
        {"kind": "section", "title": "Part I", "subtitle": "General lecture: the field before our project"},
        {
            "kind": "split",
            "title": "Why Music Style Transfer Is Hard",
            "left_title": "The problem is not just timbre",
            "left": [
                "Genre affects instrumentation, arrangement, articulation, groove feel, textural density, and production.",
                "A superficial timbre swap does not make the result sound as if it had been written in the target genre.",
                "Audio errors are unforgiving: phase issues, smeared attacks, clicks, and warble are obvious to listeners.",
            ],
            "right_title": "What a serious system must solve",
            "right": [
                "Separate content from style without obvious leakage.",
                "Represent the target genre as a usable control signal.",
                "Reconstruct realistic audio without destroying melodic identity.",
                "Maintain consistency over longer time scales.",
            ],
        },
        {
            "kind": "split",
            "title": "Content, Style, and Genre",
            "left_title": "Working definitions",
            "left": [
                "Content: melody, harmonic motion, rhythmic backbone, phrase placement.",
                "Style: timbre, articulation, groove, density, attack/decay, production texture.",
                "Genre: a broader manifold shaped by convention, culture, era, and recurring stylistic patterns.",
            ],
            "right_title": "Why the distinction matters",
            "right": [
                "If the target signal is too weak, the model only repaints surface texture.",
                "If content and style are not separated, strong edits erase the identity of the song.",
                "A good model moves the sample toward the target genre manifold while preserving recognizable content.",
            ],
        },
        {
            "kind": "formula",
            "title": "Problem Formulation",
            "formula": formulas["problem"],
            "bullets": [
                "The output should sound like a plausible target-genre performance, not a filtered source.",
                "This implies at least three axes of success: content preservation, style fidelity, and realism.",
                "Long-form systems add a fourth axis: coherence across chunk boundaries and phrases.",
            ],
        },
        {
            "kind": "split",
            "title": "Representation Levels in Music ML",
            "left_title": "Available levels",
            "left": [
                "Waveforms capture acoustics directly but are hard spaces for semantic editing.",
                "Time-frequency features such as log-mel, chroma, onsets, and beat grids expose structure more explicitly.",
                "Symbolic formats like MIDI or note-rolls help with long-range organization but are weak for timbral realism.",
            ],
            "right_title": "Design lesson",
            "right": [
                "Strong systems usually combine levels rather than choosing only one.",
                "Structure is often controlled in spectrogram or symbolic spaces, while realism is delegated to learned decoders or vocoders.",
                "Our own project follows that hybrid logic.",
            ],
        },
        {
            "kind": "table",
            "title": "Historical Evolution of the Field",
            "columns": ["Era", "Typical approach", "Core limitation"],
            "rows": [
                ["Early transfer", "Direct mapping and signal tricks", "Weak semantic control and obvious artifacts"],
                ["Disentanglement era", "Separate pitch, rhythm, timbre, or content/style", "Leakage is difficult to remove and audit"],
                ["High-fidelity generation", "GANs, codecs, vocoders, diffusion", "Control and training complexity increase sharply"],
                ["Long-form systems", "Hierarchical latents, anchoring, chunk policies", "Drift and accumulation remain hard"],
            ],
        },
        {
            "kind": "table",
            "title": "What Evaluation Must Measure",
            "columns": ["Axis", "Question", "Common proxy"],
            "rows": [
                ["Content preservation", "Did the song identity survive?", "latent similarity, pitch or MPS-style metrics"],
                ["Style fidelity", "Did it move toward the target genre?", "classifier confidence, centroid alignment"],
                ["Audio realism", "Does it sound plausible?", "MR-STFT, adversarial critics, listening"],
                ["Long-form coherence", "Does it stay consistent over time?", "boundary metrics, drift checks, human tests"],
            ],
        },
        {
            "kind": "formula",
            "title": "Tool 1: Disentanglement",
            "formula": formulas["disentangle"],
            "bullets": [
                "The content branch is rewarded for invariance and punished when style can be recovered from it.",
                "A gradient reversal layer makes style suppression explicit rather than implicit.",
                "Dual-encoder, U-Net, and sequence-analogy variants all try to separate factors, but none guarantee success without audits.",
            ],
        },
        {
            "kind": "formula",
            "title": "Tool 2: Target Style Spaces",
            "formula": formulas["target160"],
            "bullets": [
                "A target genre should be treated as a region in style space, not a single example or flat one-hot label.",
                "Good target spaces have stable centroids, useful nearest-neighbor behavior, and clear class separation.",
                "If the style space collapses, every later generator looks weaker than it really is.",
            ],
        },
        {
            "kind": "formula",
            "title": "Tool 3: Conditioning Mechanisms",
            "formula": formulas["film"],
            "bullets": [
                "FiLM allows a control vector to modulate intermediate activations throughout the network.",
                "AdaIN is useful when style should act more like timbre or texture than global structure.",
                "A good architecture decides where style and structure should interact and where they should remain separate.",
            ],
        },
        {
            "kind": "table",
            "title": "Tool 4: Generator Families",
            "columns": ["Family", "Why people use it", "Tradeoff"],
            "rows": [
                ["Waveform GANs", "Fast parallel generation", "Phase and local coherence are difficult"],
                ["Spectral GANs", "Better phase handling in structured domains", "Need strong inversion or vocoding"],
                ["Neural codec editors", "Decoder acts as a realism prior", "The translator can become too conservative"],
                ["Diffusion models", "Stable training and rich controllability", "Slow inference and long-form drift"],
            ],
        },
        {
            "kind": "formula",
            "title": "Tool 5: Diffusion and Guidance",
            "formula": formulas["diffusion"],
            "bullets": [
                "The forward process destroys structure gradually and the reverse model learns to reconstruct it.",
                "Classifier-free guidance provides a direct knob over how strongly the model should chase the condition.",
                "This is especially useful when style intensity and content preservation compete.",
            ],
        },
        {
            "kind": "formula",
            "title": "Tool 6: Long-Form Coherence",
            "formula": formulas["prefix"],
            "bullets": [
                "SDEdit-style anchoring starts from a noised source rather than pure noise so macro-structure survives.",
                "Overlap locking tries to make adjacent chunks compatible during sampling, not only after vocoding.",
                "Long-form quality becomes a systems problem involving chunk policy, re-anchoring, and stabilization.",
            ],
        },
        {
            "kind": "split",
            "title": "Field-Level Design Rules",
            "left_title": "What to do",
            "left": [
                "Measure representation quality separately from generation quality.",
                "Listen to checkpoints rather than trusting only scalar metrics.",
                "Use staged architectures when data and compute do not justify a monolithic model.",
            ],
            "right_title": "What to avoid",
            "right": [
                "Assuming a high genre score means the transfer is musically convincing.",
                "Treating long-form crossfades as a substitute for coherence constraints.",
                "Ignoring dataset-source confounds in the label design.",
            ],
        },
        {"kind": "section", "title": "Part II", "subtitle": "DGGR as a case study built from those ideas"},
        {
            "kind": "split",
            "title": "DGGR: Project Thesis",
            "left_title": "What we set out to build",
            "left": [
                "A staged genre-remastering pipeline rather than a shallow style filter.",
                "Lab 1 deconstructs the source, Lab 2 builds the target blueprint, Lab 3 reconstructs audio, Lab 4 extends to long-form coherence.",
                "Lab 5 is the planned perceptual validation layer.",
            ],
            "right_title": "Current repo-backed status",
            "right": [
                "Labs 1 through 4 are implemented in code.",
                "We have saved metrics, training histories, generated clips, and long-form diagnostics.",
                "The strongest current evidence path is the codec branch for short-form transfer.",
            ],
        },
        {"kind": "pipeline", "title": "DGGR Pipeline Overview"},
        {
            "kind": "split",
            "title": "Data Universe and Label Risks",
            "left_title": "What the workflow currently uses",
            "left": [
                "Baroque/classical renders, hip-hop material, lo-fi material, open-domain CC0 music, and speech negatives for the gate.",
                "Observed Lab 2 class buckets: baroque_classical 411, cc0_other 1200, hiphop_xtc 1200, lofi_hh_lfbb 1200.",
            ],
            "right_title": "Why this is a real risk",
            "right": [
                "Genre labels can correlate with dataset source and recording conditions.",
                "A model may learn source fingerprints instead of transferable style.",
                "This is not a side note; it changes how we interpret every later metric.",
            ],
        },
        {
            "kind": "split",
            "title": "Lab 1: Deconstruction Encoder",
            "left_title": "Architecture",
            "left": [
                "Input: 96-bin log-mel spectrogram from a fixed 5-second chunk.",
                "Backbone: 3-block Conv2D encoder with stride-2 downsampling, global pooling, and shared projection.",
                "Heads: z_content, z_style, style probe, adversarial style-from-content probe, and music gate.",
            ],
            "right_title": "Why this matters",
            "right": [
                "Every later module assumes style and content can be recombined cleanly.",
                "The training curriculum increases adversarial pressure and sharpens gate behavior over phases.",
                "Without this step, the later generator becomes a brittle edit network.",
            ],
        },
        {
            "kind": "split",
            "title": "Lab 1: Measured Outcome",
            "left_title": "Best observed audit",
            "left": [
                "Style probe accuracy: 0.9417",
                "Content leakage above baseline: 0.1083",
                "Gate ROC AUC: 0.9299",
                "All three pass the project thresholds.",
            ],
            "right_title": "Interpretation",
            "right": [
                "z_style is informative.",
                "z_content is substantially style-suppressed.",
                "The gate separates music from non-music reliably, though thresholding still trades precision for recall.",
            ],
        },
        {
            "kind": "image",
            "title": "Lab 2: Building the Target160 Style Space",
            "image": ASSETS_DIR / "lab2_tsne.png",
            "caption": "t-SNE from the current Lab 2 calibration artifacts",
            "bullets": [
                "target160 combines a 128D learned style vector with 32 mel-descriptor statistics.",
                "Per-genre centroids are computed after inlier filtering to avoid blueprint drift.",
                "The point of the space is not only separation but generator-usable geometry.",
            ],
        },
        {
            "kind": "split",
            "title": "Lab 2: Validation",
            "left_title": "Current metrics",
            "left": [
                "Linear probe accuracy: 0.8554",
                "Nearest centroid accuracy: 0.8514",
                "Silhouette (cosine): 0.4939",
                "The project threshold was silhouette >= 0.45.",
            ],
            "right_title": "What the result means",
            "right": [
                "The target space is separated enough to act as a real conditioning blueprint.",
                "This is why later improvements in generation are believable rather than accidental.",
                "Lab 2 is upstream of Lab 3 quality.",
            ],
        },
        {
            "kind": "split",
            "title": "Lab 3A: Codec-Latent Translator",
            "left_title": "Architecture",
            "left": [
                "Operate on EnCodec quantized embeddings rather than raw waveform samples.",
                "Conv1D in-projection, hidden width 256, 10 FiLM-conditioned residual blocks, then Conv1D back to 128 channels.",
                "Condition vector combines z_content, target style, and small noise.",
            ],
            "right_title": "Key decision",
            "right": [
                "The best run uses direct-output mode instead of residual mode.",
                "That removes the identity leash that was limiting style-shift magnitude.",
                "The EnCodec decoder provides a realism prior for short-form generation.",
            ],
        },
        {
            "kind": "formula",
            "title": "Lab 3A: Loss Stack and Best Run",
            "formula": formulas["codec_loss"],
            "bullets": [
                "Best saved run: run1055.",
                "MPS 0.9565, style confidence 0.8940, style accuracy 0.9492.",
                "The strongest explanation is the combination of MERT-based conditioning, direct output, and explicit content preservation pressure.",
            ],
        },
        {
            "kind": "split",
            "title": "Lab 3B: Diffusion V2",
            "left_title": "Architecture",
            "left": [
                "Input tensor: [B, 15, 80, 432] = noisy mel + chroma + onset + beat-grid channels.",
                "UNet channels [64, 128, 256, 256], two residual blocks per level, low-resolution attention, EMA shadow model.",
                "Output: [B, 1, 80, 432] velocity/noise prediction for DDIM sampling.",
            ],
            "right_title": "Condition split",
            "right": [
                "time + content -> FiLM",
                "style -> dedicated StyleAdaIN",
                "This separation is intended to keep texture control from overwhelming melody control.",
            ],
        },
        {
            "kind": "image",
            "title": "Diffusion Training Behavior",
            "image": ASSETS_DIR / "diffusion_v2_curve.png",
            "caption": "validation loss history for the current V2 run",
            "bullets": [
                "Best numeric validation loss occurs late, but the team selected epoch 6 as the best-sounding checkpoint.",
                "This is a standard generative-audio lesson: the perceptual optimum can arrive before the scalar optimum.",
                "The diffusion branch remains the more ambitious but less reliable reconstruction path.",
            ],
        },
        {
            "kind": "image",
            "title": "Lab 4: Long-Form Coherence Diagnostics",
            "image": ASSETS_DIR / "longform_comparison.png",
            "caption": "full-song coherence diagnostics over 160 seconds",
            "bullets": [
                "64 overlapping chunks over a 160-second song.",
                "Boundary mel MSE mean 0.0018347, boundary discontinuity mean 2.8691 dB.",
                "Seams are not the main failure mode; warble and static accumulation are harder problems now.",
            ],
        },
        {
            "kind": "image",
            "title": "Quantitative Summary Across Labs",
            "image": ASSETS_DIR / "metrics_dashboard.png",
            "caption": "repo-backed metric summary generated from saved artifacts",
            "bullets": [
                "Lab 1 and Lab 2 both clear their target thresholds.",
                "The codec branch is the strongest current short-form evidence path.",
                "Open gaps remain in diffusion quality, long-form texture stability, and human evaluation.",
            ],
        },
        {
            "kind": "audio",
            "title": "Demo Slide: Short-Form Codec Examples",
            "image": ASSETS_DIR / "demo_audio_grid.png",
            "caption": "spectrogram montage from the best codec run",
            "audio": [
                ("CC0 other to lo-fi", MEDIA_DIR / "codec_src1_tgt3.mp3"),
                ("hip-hop to open-domain other", MEDIA_DIR / "codec_src2_tgt1.mp3"),
                ("lo-fi to baroque", MEDIA_DIR / "codec_src3_tgt0.mp3"),
            ],
            "bullets": [
                "Listen for melody survival, target-genre plausibility, and whether the output sounds like audio rather than a synthetic artifact.",
            ],
        },
        {
            "kind": "audio",
            "title": "Demo Slide: Diffusion and Long-Form Excerpts",
            "image": ASSETS_DIR / "longform_excerpt_comparison.png",
            "caption": "source/remaster excerpt comparison plus diffusion example",
            "audio": [
                ("Diffusion V2 sample", MEDIA_DIR / "diffusion_v2_epoch6_gen0.mp3"),
                ("Source excerpt 30s to 50s", MEDIA_DIR / "longform_source_excerpt.mp3"),
                ("Remastered excerpt same segment", MEDIA_DIR / "longform_remaster_excerpt.mp3"),
            ],
            "bullets": [
                "Here the listening question changes: continuity, phrasing, and artifact buildup matter more than isolated short-form quality.",
            ],
        },
        {
            "kind": "split",
            "title": "Creative Elements and Limitations",
            "left_title": "Why the project is not generic",
            "left": [
                "Structure-first framing of genre remastering rather than shallow timbre swapping.",
                "A staged architecture with explicit deconstruction, calibration, reconstruction, and coherence stages.",
                "Direct-output codec translation, MERT-based conditioning, and prefix-lock sampling controls.",
            ],
            "right_title": "Where the work remains",
            "right": [
                "Dataset-source leakage is still a real risk.",
                "Diffusion quality still trails codec transfer on short-form reliability.",
                "Long-form audio still accumulates warble and static under stronger edits.",
            ],
        },
        {
            "kind": "split",
            "title": "Lab 5 and Conclusion",
            "left_title": "What Lab 5 must answer",
            "left": [
                "Can listeners identify the target genre?",
                "Can they still recognize the source melody?",
                "Do they prefer our remaster to a simpler filter-style baseline?",
            ],
            "right_title": "Bottom line",
            "right": [
                "The repo already supports a credible architecture story and measurable intermediate success.",
                "The remaining question is no longer whether the system runs.",
                "The real question is whether people actually hear the intended effect convincingly.",
            ],
        },
    ]


def add_footer(slide, idx: int, total: int) -> None:
    tb = slide.shapes.add_textbox(Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.22))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{idx}/{total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = rgb("muted")


def set_text_style(run, size: int, color_name: str = "ink", bold: bool = False) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color_name)


def setup_text_frame(tf, wrap: bool = True, margin: int = 0) -> None:
    tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(margin)
    tf.margin_right = Pt(margin)
    tf.margin_top = Pt(margin)
    tf.margin_bottom = Pt(margin)


def add_bullet_paragraph(tf, text: str, size: int = 15, color_name: str = "ink") -> None:
    p = tf.add_paragraph()
    p.text = f"• {text}"
    p.space_after = Pt(4)
    p.font.size = Pt(size)
    p.font.color.rgb = rgb(color_name)


def add_body_textbox(slide, left: float, top: float, width: float, height: float, fill: str = "paper"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb("slate")
    shape.line.width = Pt(1.0)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    return tf


def estimate_lines(text: str, chars_per_line: int) -> int:
    text = text.strip()
    if not text:
        return 1
    return max(1, (len(text) + chars_per_line - 1) // chars_per_line)


def estimate_card_height_pptx(title: str, bullets: list[str], width: float, body_font: int = 15) -> float:
    chars_per_line = max(24, int(width * 6.2))
    title_lines = estimate_lines(title, chars_per_line - 4)
    bullet_lines = sum(estimate_lines(bullet, chars_per_line - 2) for bullet in bullets)
    line_unit = 0.28 if body_font >= 15 else 0.24
    base = 0.78 + 0.22 * max(0, title_lines - 1)
    height = base + bullet_lines * line_unit + max(0, len(bullets) - 1) * 0.09
    return max(1.8, min(height, 4.9))


def estimate_card_height_pdf(title: str, bullets: list[str], width: float, body_font: int = 11) -> float:
    chars_per_line = max(36, int(width / 7.6))
    title_lines = estimate_lines(title, chars_per_line - 4)
    bullet_lines = sum(estimate_lines(bullet, chars_per_line - 2) for bullet in bullets)
    line_unit = 16 if body_font >= 11 else 14
    base = 34 + 8 * max(0, title_lines - 1)
    height = base + bullet_lines * line_unit + max(0, len(bullets) - 1) * 5
    return max(105, min(height, 360))


def collect_unique_audio_files(slides: list[dict]) -> list[Path]:
    seen: set[Path] = set()
    ordered: list[Path] = []
    for spec in slides:
        if spec.get("kind") != "audio":
            continue
        for _, path in spec["audio"]:
            resolved = path.resolve()
            if resolved in seen:
                continue
            seen.add(resolved)
            ordered.append(resolved)
    return ordered


def render_title_pptx(prs: Presentation, spec: dict, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("sand")
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.05), Inches(5.35), Inches(2.7), Inches(2.7))
    accent.fill.solid(); accent.fill.fore_color.rgb = rgb("mist"); accent.line.fill.background()
    accent2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.45), Inches(0.0), Inches(2.2), Inches(2.2))
    accent2.fill.solid(); accent2.fill.fore_color.rgb = rgb("peach"); accent2.line.fill.background()
    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.72), Inches(8.15), Inches(1.82))
    setup_text_frame(title.text_frame)
    p = title.text_frame.paragraphs[0]
    run = p.add_run(); run.text = spec["title"]; set_text_style(run, 30, "ink", True)
    sub = slide.shapes.add_textbox(Inches(0.77), Inches(2.02), Inches(6.1), Inches(0.55))
    setup_text_frame(sub.text_frame)
    p = sub.text_frame.paragraphs[0]
    run = p.add_run(); run.text = spec["subtitle"]; set_text_style(run, 18, "muted", False)
    goal_height = estimate_card_height_pptx("Goal", spec["body"], 4.15, 16)
    tf = add_body_textbox(slide, 7.8, 0.92, 4.15, goal_height, "paper")
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "Goal"; set_text_style(r, 17, "teal", True)
    for bullet in spec["body"]:
        add_bullet_paragraph(tf, bullet, 16, "ink")
    names = slide.shapes.add_textbox(Inches(0.77), Inches(2.75), Inches(6.8), Inches(0.6))
    setup_text_frame(names.text_frame)
    p = names.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Sahara Kaul  •  Kelsey Pattison  •  Ahmed Sajid"; set_text_style(r, 15, "ink", False)
    c = slide.shapes.add_textbox(Inches(0.77), Inches(3.17), Inches(6.1), Inches(0.4))
    setup_text_frame(c.text_frame)
    p = c.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "CMPUT 414, Winter 2026"; set_text_style(r, 13, "muted", False)
    add_footer(slide, idx, total)


def render_section_pptx(prs: Presentation, spec: dict, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("night")
    title = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(10.5), Inches(1.2))
    setup_text_frame(title.text_frame)
    p = title.text_frame.paragraphs[0]
    r = p.add_run(); r.text = spec["title"]; set_text_style(r, 30, "paper", True)
    sub = slide.shapes.add_textbox(Inches(0.92), Inches(3.0), Inches(10.6), Inches(0.8))
    setup_text_frame(sub.text_frame)
    p = sub.text_frame.paragraphs[0]
    r = p.add_run(); r.text = spec["subtitle"]; set_text_style(r, 18, "mist", False)
    add_footer(slide, idx, total)


def render_split_pptx(prs: Presentation, spec: dict, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb("sand")
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(11.8), Inches(0.58))
    setup_text_frame(title.text_frame)
    p = title.text_frame.paragraphs[0]; r = p.add_run(); r.text = spec["title"]; set_text_style(r, 26, "ink", True)
    left_height = estimate_card_height_pptx(spec["left_title"], spec["left"], 5.55, 16)
    right_height = estimate_card_height_pptx(spec["right_title"], spec["right"], 5.55, 16)
    left = add_body_textbox(slide, 0.55, 1.0, 5.55, left_height, "paper")
    p = left.paragraphs[0]; r = p.add_run(); r.text = spec["left_title"]; set_text_style(r, 17, "teal", True)
    for bullet in spec["left"]:
        add_bullet_paragraph(left, bullet, 16, "ink")
    right = add_body_textbox(slide, 6.25, 1.0, 5.55, right_height, "mist")
    p = right.paragraphs[0]; r = p.add_run(); r.text = spec["right_title"]; set_text_style(r, 17, "gold", True)
    for bullet in spec["right"]:
        add_bullet_paragraph(right, bullet, 16, "ink")
    add_footer(slide, idx, total)


def render_formula_pptx(prs: Presentation, spec: dict, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb("sand")
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(11.8), Inches(0.58))
    setup_text_frame(title.text_frame)
    p = title.text_frame.paragraphs[0]; r = p.add_run(); r.text = spec["title"]; set_text_style(r, 25, "ink", True)
    slide.shapes.add_picture(str(spec["formula"]), Inches(0.6), Inches(1.15), width=Inches(7.2))
    card_height = estimate_card_height_pptx("Key points", spec["bullets"], 3.95, 16)
    tf = add_body_textbox(slide, 7.7, 1.1, 3.95, card_height, "paper")
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "Key points"; set_text_style(r, 17, "teal", True)
    for bullet in spec["bullets"]:
        add_bullet_paragraph(tf, bullet, 16, "ink")
    add_footer(slide, idx, total)


def render_table_pptx(prs: Presentation, spec: dict, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb("sand")
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(11.8), Inches(0.58))
    setup_text_frame(title.text_frame)
    p = title.text_frame.paragraphs[0]; r = p.add_run(); r.text = spec["title"]; set_text_style(r, 24, "ink", True)
    rows = len(spec["rows"]) + 1
    cols = len(spec["columns"])
    table = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.15), Inches(11.1), Inches(5.8)).table
    widths = [2.1, 4.0, 5.0]
    for i, w in enumerate(widths[:cols]):
        table.columns[i].width = Inches(w)
    for i, col in enumerate(spec["columns"]):
        cell = table.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb("teal")
        cell.text = col
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = rgb("paper")
    for r_idx, row in enumerate(spec["rows"], start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb("paper" if r_idx % 2 else "mist")
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12); p.font.color.rgb = rgb("ink")
    add_footer(slide, idx, total)


def render_image_pptx(prs: Presentation, spec: dict, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb("sand")
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(11.8), Inches(0.58))
    setup_text_frame(title.text_frame)
    p = title.text_frame.paragraphs[0]; r = p.add_run(); r.text = spec["title"]; set_text_style(r, 25, "ink", True)
    slide.shapes.add_picture(str(spec["image"]), Inches(0.55), Inches(1.05), width=Inches(6.9), height=Inches(4.95))
    card_height = estimate_card_height_pptx("Why this slide matters", spec["bullets"], 4.0, 15)
    tf = add_body_textbox(slide, 7.55, 1.05, 4.0, card_height, "paper")
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "Why this slide matters"; set_text_style(r, 17, "teal", True)
    for bullet in spec["bullets"]:
        add_bullet_paragraph(tf, bullet, 15, "ink")
    cap = slide.shapes.add_textbox(Inches(0.57), Inches(6.08), Inches(6.9), Inches(0.35))
    setup_text_frame(cap.text_frame)
    p = cap.text_frame.paragraphs[0]; r = p.add_run(); r.text = spec["caption"]; set_text_style(r, 10, "muted", False)
    add_footer(slide, idx, total)


def render_pipeline_pptx(prs: Presentation, spec: dict, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb("sand")
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.2), Inches(0.5))
    setup_text_frame(title.text_frame)
    p = title.text_frame.paragraphs[0]; r = p.add_run(); r.text = spec["title"]; set_text_style(r, 24, "ink", True)
    nodes = [
        (0.7, 2.4, 1.7, 1.0, "peach", "Input", "song or 5-second chunk"),
        (2.7, 2.4, 1.9, 1.0, "mist", "Lab 1", "deconstruction\nz_content / z_style"),
        (4.9, 2.4, 1.9, 1.0, "paper", "Lab 2", "target160\ncentroids"),
        (7.2, 1.6, 1.9, 1.0, "slate", "Lab 3A", "codec translator\nEnCodec + FiLM"),
        (7.2, 3.2, 1.9, 1.0, "slate", "Lab 3B", "diffusion V2\nmel + BigVGAN"),
        (9.6, 2.4, 1.9, 1.0, "peach", "Lab 4", "long-form coherence\nprefix lock"),
    ]
    for left, top, width, height, fill, label, desc in nodes:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill); shape.line.color.rgb = rgb("slate")
        tf = shape.text_frame
        tf.margin_left = Pt(8); tf.margin_top = Pt(8)
        p = tf.paragraphs[0]; r = p.add_run(); r.text = label; set_text_style(r, 15, "ink", True)
        p = tf.add_paragraph(); p.text = desc; p.font.size = Pt(11); p.font.color.rgb = rgb("ink")
    for x1, y1, x2, y2 in [(2.4, 2.9, 2.7, 2.9), (4.6, 2.9, 4.9, 2.9), (6.8, 2.9, 7.2, 2.1), (6.8, 2.9, 7.2, 3.7), (9.1, 2.1, 9.6, 2.9), (9.1, 3.7, 9.6, 2.9)]:
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = rgb("gold")
        line.line.width = Pt(2.5)
    tf = add_body_textbox(slide, 0.8, 5.1, 10.8, 1.2, "paper")
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "Thesis"; set_text_style(r, 15, "teal", True)
    p = tf.add_paragraph(); p.text = "The system is staged because each lab solves one field-level problem explicitly instead of hiding everything inside a single end-to-end mapping."; p.font.size = Pt(14); p.font.color.rgb = rgb("ink")
    add_footer(slide, idx, total)


def render_audio_pptx(prs: Presentation, spec: dict, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb("sand")
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.2), Inches(0.5))
    setup_text_frame(title.text_frame)
    p = title.text_frame.paragraphs[0]; r = p.add_run(); r.text = spec["title"]; set_text_style(r, 25, "ink", True)
    slide.shapes.add_picture(str(spec["image"]), Inches(0.55), Inches(1.05), width=Inches(6.85), height=Inches(4.85))
    card_bullets = [f"{label}: {path.name}" for label, path in spec["audio"]] + spec["bullets"]
    card_height = estimate_card_height_pptx("Clip inventory", card_bullets, 4.05, 14)
    tf = add_body_textbox(slide, 7.5, 1.05, 4.05, card_height, "paper")
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "Clip inventory"; set_text_style(r, 17, "teal", True)
    icon_left = 7.73
    label_left = 8.33
    row_top = 1.6
    row_gap = 0.72
    for clip_idx, (label, path) in enumerate(spec["audio"]):
        top = row_top + clip_idx * row_gap
        slide.shapes.add_movie(str(path), Inches(icon_left), Inches(top), Inches(0.48), Inches(0.48), mime_type="audio/mpeg")
        label_box = slide.shapes.add_textbox(Inches(label_left), Inches(top - 0.02), Inches(3.0), Inches(0.5))
        setup_text_frame(label_box.text_frame)
        p = label_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = label
        set_text_style(run, 14, "ink", True)
        p.space_after = Pt(0)
        p = label_box.text_frame.add_paragraph()
        run = p.add_run()
        run.text = path.name
        set_text_style(run, 10, "muted", False)
    note_box = slide.shapes.add_textbox(Inches(7.73), Inches(row_top + len(spec["audio"]) * row_gap + 0.12), Inches(3.5), Inches(1.0))
    setup_text_frame(note_box.text_frame)
    p = note_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = spec["bullets"][0]
    set_text_style(run, 13, "muted", False)
    cap = slide.shapes.add_textbox(Inches(0.57), Inches(5.98), Inches(6.9), Inches(0.35))
    setup_text_frame(cap.text_frame)
    p = cap.text_frame.paragraphs[0]; r = p.add_run(); r.text = spec["caption"]; set_text_style(r, 10, "muted", False)
    add_footer(slide, idx, total)


def build_pptx(slides: list[dict]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    renderers = {
        "title": render_title_pptx,
        "section": render_section_pptx,
        "split": render_split_pptx,
        "formula": render_formula_pptx,
        "table": render_table_pptx,
        "image": render_image_pptx,
        "pipeline": render_pipeline_pptx,
        "audio": render_audio_pptx,
    }
    total = len(slides)
    for idx, spec in enumerate(slides, start=1):
        renderers[spec["kind"]](prs, spec, idx, total)
    try:
        prs.save(PPTX_OUT)
    except PermissionError:
        fallback = next_available_path(PPTX_OUT)
        prs.save(fallback)
        print(fallback)


PDF_WIDTH = 13.333 * 72
PDF_HEIGHT = 7.5 * 72
styles = getSampleStyleSheet()
TITLE_STYLE = ParagraphStyle("Title", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=22, textColor=pdf_color("ink"), leading=24)
HEAD_STYLE = ParagraphStyle("Head", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=14, textColor=pdf_color("teal"), leading=16)
BODY_STYLE = ParagraphStyle("Body", parent=styles["BodyText"], fontName="Helvetica", fontSize=11, textColor=pdf_color("ink"), leading=14)
SMALL_STYLE = ParagraphStyle("Small", parent=styles["BodyText"], fontName="Helvetica", fontSize=9, textColor=pdf_color("muted"), leading=11)


def draw_paragraph(c: canvas.Canvas, text: str, style: ParagraphStyle, x: float, y: float, width: float) -> float:
    para = Paragraph(text, style)
    _, h = para.wrap(width, PDF_HEIGHT)
    para.drawOn(c, x, y - h)
    return h


def draw_card(c: canvas.Canvas, x: float, y: float, w: float, h: float, fill_name: str) -> None:
    c.setFillColor(pdf_color(fill_name))
    c.setStrokeColor(pdf_color("slate"))
    c.roundRect(x, y - h, w, h, 12, fill=1, stroke=1)


def render_title_pdf(c: canvas.Canvas, spec: dict) -> None:
    c.setFillColor(pdf_color("sand")); c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)
    c.setFillColor(pdf_color("mist")); c.circle(90, 80, 90, fill=1, stroke=0)
    c.setFillColor(pdf_color("peach")); c.circle(PDF_WIDTH - 80, PDF_HEIGHT - 60, 80, fill=1, stroke=0)
    draw_paragraph(c, spec["title"], ParagraphStyle("big", parent=TITLE_STYLE, fontSize=24, leading=28), 45, PDF_HEIGHT - 48, 490)
    draw_paragraph(c, spec["subtitle"], ParagraphStyle("sub", parent=BODY_STYLE, fontSize=14, textColor=pdf_color("muted")), 45, PDF_HEIGHT - 108, 400)
    goal_height = estimate_card_height_pdf("Goal", spec["body"], 340, 11)
    draw_card(c, 572, PDF_HEIGHT - 70, 340, goal_height, "paper")
    draw_paragraph(c, "<b>Goal</b>", HEAD_STYLE, 588, PDF_HEIGHT - 85, 300)
    y = PDF_HEIGHT - 112
    for bullet in spec["body"]:
        h = draw_paragraph(c, f"&bull; {bullet}", BODY_STYLE, 590, y, 300)
        y -= h + 4
    draw_paragraph(c, "Sahara Kaul  •  Kelsey Pattison  •  Ahmed Sajid", BODY_STYLE, 45, PDF_HEIGHT - 155, 420)
    draw_paragraph(c, "CMPUT 414, Winter 2026", SMALL_STYLE, 45, PDF_HEIGHT - 177, 250)


def render_section_pdf(c: canvas.Canvas, spec: dict) -> None:
    c.setFillColor(pdf_color("night")); c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)
    white_title = ParagraphStyle("white_title", parent=TITLE_STYLE, fontSize=28, textColor=Color(1, 1, 1))
    white_sub = ParagraphStyle("white_sub", parent=BODY_STYLE, fontSize=16, textColor=Color(0.92, 0.96, 0.95))
    draw_paragraph(c, spec["title"], white_title, 60, PDF_HEIGHT - 180, 250)
    draw_paragraph(c, spec["subtitle"], white_sub, 60, PDF_HEIGHT - 230, 620)


def render_split_pdf(c: canvas.Canvas, spec: dict) -> None:
    c.setFillColor(pdf_color("sand")); c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)
    draw_paragraph(c, spec["title"], TITLE_STYLE, 45, PDF_HEIGHT - 28, 700)
    left_height = estimate_card_height_pdf(spec["left_title"], spec["left"], 410, 11)
    right_height = estimate_card_height_pdf(spec["right_title"], spec["right"], 410, 11)
    draw_card(c, 45, PDF_HEIGHT - 78, 410, left_height, "paper")
    draw_card(c, 485, PDF_HEIGHT - 78, 410, right_height, "mist")
    draw_paragraph(c, f"<b>{spec['left_title']}</b>", HEAD_STYLE, 60, PDF_HEIGHT - 96, 370)
    y = PDF_HEIGHT - 122
    for bullet in spec["left"]:
        h = draw_paragraph(c, f"&bull; {bullet}", BODY_STYLE, 60, y, 370)
        y -= h + 4
    draw_paragraph(c, f"<b>{spec['right_title']}</b>", ParagraphStyle("goldh", parent=HEAD_STYLE, textColor=pdf_color("gold")), 500, PDF_HEIGHT - 96, 370)
    y = PDF_HEIGHT - 122
    for bullet in spec["right"]:
        h = draw_paragraph(c, f"&bull; {bullet}", BODY_STYLE, 500, y, 370)
        y -= h + 4


def render_formula_pdf(c: canvas.Canvas, spec: dict) -> None:
    c.setFillColor(pdf_color("sand")); c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)
    draw_paragraph(c, spec["title"], TITLE_STYLE, 45, PDF_HEIGHT - 28, 700)
    c.drawImage(ImageReader(str(spec["formula"])), 45, PDF_HEIGHT - 250, width=470, height=140, preserveAspectRatio=True, mask="auto")
    card_height = estimate_card_height_pdf("Key points", spec["bullets"], 350, 11)
    draw_card(c, 545, PDF_HEIGHT - 86, 350, card_height, "paper")
    draw_paragraph(c, "<b>Key points</b>", HEAD_STYLE, 560, PDF_HEIGHT - 104, 310)
    y = PDF_HEIGHT - 130
    for bullet in spec["bullets"]:
        h = draw_paragraph(c, f"&bull; {bullet}", BODY_STYLE, 560, y, 300)
        y -= h + 4


def render_table_pdf(c: canvas.Canvas, spec: dict) -> None:
    c.setFillColor(pdf_color("sand")); c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)
    draw_paragraph(c, spec["title"], TITLE_STYLE, 45, PDF_HEIGHT - 28, 700)
    x0, y0 = 45, PDF_HEIGHT - 92
    row_h = 64 if len(spec["rows"]) <= 4 else 56
    col_w = [160, 270, 430]
    for i, col in enumerate(spec["columns"]):
        c.setFillColor(pdf_color("teal")); c.setStrokeColor(pdf_color("teal"))
        c.rect(x0 + sum(col_w[:i]), y0 - row_h, col_w[i], row_h, fill=1, stroke=1)
        draw_paragraph(c, f"<font color='white'><b>{col}</b></font>", ParagraphStyle("th", parent=BODY_STYLE, fontSize=11, textColor=Color(1, 1, 1)), x0 + sum(col_w[:i]) + 8, y0 - 10, col_w[i] - 16)
    for ridx, row in enumerate(spec["rows"], start=1):
        yy = y0 - row_h * (ridx + 1)
        fill = "paper" if ridx % 2 else "mist"
        for cidx, val in enumerate(row):
            xx = x0 + sum(col_w[:cidx])
            c.setFillColor(pdf_color(fill)); c.setStrokeColor(pdf_color("slate"))
            c.rect(xx, yy, col_w[cidx], row_h, fill=1, stroke=1)
            draw_paragraph(c, val, ParagraphStyle("cell", parent=BODY_STYLE, fontSize=10), xx + 8, yy + row_h - 8, col_w[cidx] - 16)


def render_image_pdf(c: canvas.Canvas, spec: dict) -> None:
    c.setFillColor(pdf_color("sand")); c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)
    draw_paragraph(c, spec["title"], TITLE_STYLE, 45, PDF_HEIGHT - 28, 700)
    c.drawImage(ImageReader(str(spec["image"])), 45, PDF_HEIGHT - 430, width=500, height=300, preserveAspectRatio=True, mask="auto")
    draw_paragraph(c, spec["caption"], SMALL_STYLE, 47, PDF_HEIGHT - 437, 500)
    card_height = estimate_card_height_pdf("Why this slide matters", spec["bullets"], 330, 11)
    draw_card(c, 565, PDF_HEIGHT - 86, 330, card_height, "paper")
    draw_paragraph(c, "<b>Why this slide matters</b>", HEAD_STYLE, 580, PDF_HEIGHT - 104, 295)
    y = PDF_HEIGHT - 130
    for bullet in spec["bullets"]:
        h = draw_paragraph(c, f"&bull; {bullet}", BODY_STYLE, 580, y, 295)
        y -= h + 4


def render_pipeline_pdf(c: canvas.Canvas, spec: dict) -> None:
    c.setFillColor(pdf_color("sand")); c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)
    draw_paragraph(c, spec["title"], TITLE_STYLE, 45, PDF_HEIGHT - 28, 700)
    nodes = [
        (60, PDF_HEIGHT - 260, 110, 55, "peach", "Input", "song / chunk"),
        (190, PDF_HEIGHT - 260, 130, 55, "mist", "Lab 1", "deconstruction"),
        (345, PDF_HEIGHT - 260, 130, 55, "paper", "Lab 2", "target160"),
        (520, PDF_HEIGHT - 190, 140, 55, "slate", "Lab 3A", "codec"),
        (520, PDF_HEIGHT - 330, 140, 55, "slate", "Lab 3B", "diffusion"),
        (720, PDF_HEIGHT - 260, 140, 55, "peach", "Lab 4", "coherence"),
    ]
    for x, y, w, h, fill, t1, t2 in nodes:
        c.setFillColor(pdf_color(fill)); c.setStrokeColor(pdf_color("slate"))
        c.roundRect(x, y, w, h, 10, fill=1, stroke=1)
        draw_paragraph(c, f"<b>{t1}</b><br/>{t2}", ParagraphStyle("node", parent=BODY_STYLE, fontSize=11), x + 8, y + h - 8, w - 16)
    c.setStrokeColor(pdf_color("gold")); c.setLineWidth(2)
    for x1, y1, x2, y2 in [(170, PDF_HEIGHT - 233, 190, PDF_HEIGHT - 233), (320, PDF_HEIGHT - 233, 345, PDF_HEIGHT - 233), (475, PDF_HEIGHT - 233, 520, PDF_HEIGHT - 163), (475, PDF_HEIGHT - 233, 520, PDF_HEIGHT - 303), (660, PDF_HEIGHT - 163, 720, PDF_HEIGHT - 233), (660, PDF_HEIGHT - 303, 720, PDF_HEIGHT - 233)]:
        c.line(x1, y1, x2, y2)
    draw_card(c, 60, PDF_HEIGHT - 405, 800, 82, "paper")
    draw_paragraph(c, "<b>Thesis</b><br/>The system is staged because each lab solves one field-level problem explicitly instead of hiding everything inside one monolithic mapping.", BODY_STYLE, 78, PDF_HEIGHT - 337, 760)


def render_audio_pdf(c: canvas.Canvas, spec: dict) -> None:
    c.setFillColor(pdf_color("sand")); c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)
    draw_paragraph(c, spec["title"], TITLE_STYLE, 45, PDF_HEIGHT - 28, 700)
    c.drawImage(ImageReader(str(spec["image"])), 45, PDF_HEIGHT - 410, width=500, height=280, preserveAspectRatio=True, mask="auto")
    draw_paragraph(c, spec["caption"], SMALL_STYLE, 47, PDF_HEIGHT - 417, 500)
    card_bullets = [f"{label}: {path.name}" for label, path in spec["audio"]] + spec["bullets"]
    card_height = estimate_card_height_pdf("Clip inventory", card_bullets, 330, 10)
    draw_card(c, 565, PDF_HEIGHT - 86, 330, card_height, "paper")
    draw_paragraph(c, "<b>Clip inventory</b>", HEAD_STYLE, 580, PDF_HEIGHT - 104, 290)
    y = PDF_HEIGHT - 132
    for label, path in spec["audio"]:
        h = draw_paragraph(c, f"<b>{label}</b><br/><font face='Courier'>{path.name}</font>", BODY_STYLE, 580, y, 290)
        y -= h + 6
    draw_paragraph(c, spec["bullets"][0], SMALL_STYLE, 580, y - 4, 290)


def build_pdf(slides: list[dict]) -> None:
    c = canvas.Canvas(str(PDF_OUT), pagesize=(PDF_WIDTH, PDF_HEIGHT))
    renderers = {
        "title": render_title_pdf,
        "section": render_section_pdf,
        "split": render_split_pdf,
        "formula": render_formula_pdf,
        "table": render_table_pdf,
        "image": render_image_pdf,
        "pipeline": render_pipeline_pdf,
        "audio": render_audio_pdf,
    }
    total = len(slides)
    for idx, spec in enumerate(slides, start=1):
        renderers[spec["kind"]](c, spec)
        c.setFillColor(pdf_color("muted"))
        c.setFont("Helvetica", 9)
        c.drawRightString(PDF_WIDTH - 30, 18, f"{idx}/{total}")
        c.showPage()
    c.save()


def attach_pdf_audio(slides: list[dict]) -> None:
    doc = fitz.open(str(PDF_OUT))
    embedded = set(doc.embfile_names())
    for path in collect_unique_audio_files(slides):
        if path.name not in embedded:
            doc.embfile_add(
                path.name,
                path.read_bytes(),
                filename=path.name,
                ufilename=path.name,
                desc=f"DGGR demo audio clip: {path.name}",
            )
    for page_idx, spec in enumerate(slides):
        if spec.get("kind") != "audio":
            continue
        page = doc[page_idx]
        for clip_idx, (label, path) in enumerate(spec["audio"]):
            point = fitz.Point(580, 108 + clip_idx * 52)
            page.add_file_annot(
                point,
                path.read_bytes(),
                filename=path.name,
                ufilename=path.name,
                desc=label,
                icon="PushPin",
            )
    tmp_out = PDF_OUT.with_name(f"{PDF_OUT.stem}_with_audio{PDF_OUT.suffix}")
    doc.save(tmp_out, garbage=4, deflate=True)
    doc.close()
    tmp_out.replace(PDF_OUT)


def main() -> None:
    ensure_dirs()
    formulas = formula_assets()
    slides = slide_specs(formulas)
    build_pptx(slides)
    build_pdf(slides)
    attach_pdf_audio(slides)
    print(PPTX_OUT)
    print(PDF_OUT)


if __name__ == "__main__":
    main()
